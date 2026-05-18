from __future__ import annotations

import hashlib
import mimetypes
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol
import wave

from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation

from app.core.config import settings
from app.core import session_assets as session_asset_utils
from app.providers import SpeechRecognitionRequest
from app.schemas.assets import SessionAssetSummary
import app.services.asset_service as asset_service
import app.services.speech_recognition_service as speech_recognition_service


class AssetUploadError(ValueError):
    pass


class AsyncUploadLike(Protocol):
    filename: str | None
    content_type: str | None

    async def read(self, size: int = -1) -> bytes: ...


@dataclass
class ParsedChunk:
    content: str
    page_number: int | None = None
    sheet_name: str | None = None
    slide_number: int | None = None
    section_path: str | None = None
    start_ms: int | None = None
    end_ms: int | None = None
    speaker: str | None = None
    frame_index: int | None = None
    frame_timestamp_ms: int | None = None
    summary: str | None = None


_AUDIO_SUFFIXES = {".wav", ".mp3", ".m4a", ".aac", ".ogg", ".flac"}
_VIDEO_SUFFIXES = {".mp4", ".mov", ".m4v", ".webm", ".mkv"}


def detect_asset_kind(filename: str, mime_type: str) -> str:
    suffix = Path(filename).suffix.lower()
    lowered_mime = (mime_type or "").lower()
    if lowered_mime.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}:
        return "image"
    if lowered_mime.startswith("audio/") or suffix in _AUDIO_SUFFIXES:
        return "audio"
    if lowered_mime.startswith("video/") or suffix in _VIDEO_SUFFIXES:
        return "video"
    if lowered_mime == "application/pdf" or suffix == ".pdf":
        return "pdf"
    if suffix == ".docx":
        return "docx"
    if suffix == ".xlsx":
        return "xlsx"
    if suffix == ".pptx":
        return "pptx"
    return "other"


def normalize_mime_type(filename: str, mime_type: str | None) -> str:
    if mime_type and mime_type.strip():
        return mime_type.strip().lower()
    guessed, _ = mimetypes.guess_type(filename)
    return (guessed or "application/octet-stream").lower()


def validate_upload(filename: str, mime_type: str, size_bytes: int) -> str:
    if not filename.strip():
        raise AssetUploadError("Uploaded file is missing a filename.")
    if size_bytes <= 0:
        raise AssetUploadError("Uploaded file is empty.")
    if size_bytes > settings.jarvis_asset_max_file_bytes:
        raise AssetUploadError("Uploaded file exceeds the maximum allowed size.")
    kind = detect_asset_kind(filename, mime_type)
    if kind == "image" and size_bytes > settings.jarvis_asset_max_image_bytes:
        raise AssetUploadError("Uploaded image exceeds the maximum allowed size.")
    if kind == "other":
        raise AssetUploadError("Unsupported file type. Jarvis currently accepts images, audio, video, PDF, DOCX, XLSX, and PPTX.")
    return kind


def stage_uploaded_asset(session_id: str, filename: str, mime_type: str | None, data: bytes) -> SessionAssetSummary:
    normalized_mime = normalize_mime_type(filename, mime_type)
    kind = validate_upload(filename, normalized_mime, len(data))
    sha256 = hashlib.sha256(data).hexdigest()
    asset = asset_service.create_asset_record(
        session_id,
        kind=kind,
        mime_type=normalized_mime,
        filename=filename,
        size_bytes=len(data),
        sha256=sha256,
        status="uploaded",
    )
    session_asset_utils.ensure_asset_dirs(session_id, asset.id)
    original_path = Path(asset.storage_path)
    original_path.write_bytes(data)
    updated = asset_service.update_asset_record(asset.id, storage_path=original_path.as_posix(), sha256=sha256)
    return updated or asset


async def stage_uploaded_asset_stream(
    session_id: str,
    upload: AsyncUploadLike,
    *,
    chunk_size: int = 1024 * 1024,
) -> SessionAssetSummary:
    filename = upload.filename or ""
    normalized_mime = normalize_mime_type(filename, upload.content_type)
    kind = detect_asset_kind(filename, normalized_mime)
    if kind == "other":
        raise AssetUploadError("Unsupported file type. Jarvis currently accepts images, audio, video, PDF, DOCX, XLSX, and PPTX.")

    asset_id = session_asset_utils.new_asset_id()
    session_asset_utils.ensure_asset_dirs(session_id, asset_id)
    original_path = session_asset_utils.allocate_original_path(session_id, asset_id, filename)

    sha256_hasher = hashlib.sha256()
    size_bytes = 0
    try:
        with original_path.open("wb") as handle:
            while True:
                chunk = await upload.read(chunk_size)
                if not chunk:
                    break
                handle.write(chunk)
                sha256_hasher.update(chunk)
                size_bytes += len(chunk)

        validate_upload(filename, normalized_mime, size_bytes)
        asset = asset_service.create_asset_record(
            session_id,
            asset_id=asset_id,
            kind=kind,
            mime_type=normalized_mime,
            filename=filename,
            size_bytes=size_bytes,
            sha256=sha256_hasher.hexdigest(),
            status="uploaded",
            storage_path=original_path.as_posix(),
        )
        return asset
    except Exception:
        try:
            original_path.unlink(missing_ok=True)
        except Exception:
            pass
        raise


def ingest_asset(asset_id: str) -> SessionAssetSummary:
    asset = asset_service.get_asset(asset_id)
    if asset is None:
        raise AssetUploadError(f"Unknown asset {asset_id}")
    asset_service.delete_asset_chunks(asset.id)
    asset_service.update_asset_record(asset.id, status="processing", error_message="")
    refreshed = asset_service.get_asset(asset.id)
    if refreshed is None:
        raise AssetUploadError(f"Unknown asset {asset_id}")

    try:
        if refreshed.kind == "image":
            preview_path = _ingest_image(refreshed)
            return asset_service.update_asset_record(
                refreshed.id,
                status="ready",
                preview_path=preview_path,
                error_message="",
            ) or refreshed

        if refreshed.kind == "audio":
            metadata_json = {**refreshed.metadata_json, **_ingest_audio(refreshed)}
            return asset_service.update_asset_record(
                refreshed.id,
                status="ready",
                error_message="",
                metadata_json=metadata_json,
            ) or refreshed

        if refreshed.kind == "video":
            metadata_json = {**refreshed.metadata_json, **_ingest_video(refreshed)}
            return asset_service.update_asset_record(
                refreshed.id,
                status="ready",
                error_message="",
                metadata_json=metadata_json,
            ) or refreshed

        chunks = _extract_document_chunks(refreshed)
        extracted_text = "\n\n".join(chunk.content for chunk in chunks).strip()
        text_path = session_asset_utils.allocate_extracted_text_path(refreshed.session_id, refreshed.id)
        text_path.write_text(extracted_text)
        for index, chunk in enumerate(chunks):
            asset_service.create_asset_chunk(
                refreshed.id,
                chunk_index=index,
                content=chunk.content,
                page_number=chunk.page_number,
                sheet_name=chunk.sheet_name,
                slide_number=chunk.slide_number,
                section_path=chunk.section_path,
                start_ms=chunk.start_ms,
                end_ms=chunk.end_ms,
                speaker=chunk.speaker,
                frame_index=chunk.frame_index,
                frame_timestamp_ms=chunk.frame_timestamp_ms,
                summary=chunk.summary,
            )
        return asset_service.update_asset_record(
            refreshed.id,
            status="ready",
            error_message="",
        ) or refreshed
    except Exception as exc:
        failed = asset_service.update_asset_record(
            refreshed.id,
            status="failed",
            error_message=str(exc),
        )
        if failed is None:
            raise
        return failed


def _ingest_image(asset: SessionAssetSummary) -> str:
    original_path = Path(asset.storage_path)
    preview_path = session_asset_utils.allocate_preview_path(asset.session_id, asset.id)
    with Image.open(original_path) as image:
        preview = image.convert("RGB")
        preview.thumbnail((320, 320))
        preview.save(preview_path, format="PNG")
    return preview_path.as_posix()


def _ingest_audio(asset: SessionAssetSummary) -> dict[str, object]:
    original_path = Path(asset.storage_path)
    metadata = _probe_media_metadata(original_path, kind="audio")
    transcript_status, transcript_chunks = _transcribe_audio_asset(asset, metadata)
    metadata["transcript_status"] = transcript_status
    for index, chunk in enumerate(transcript_chunks):
        asset_service.create_asset_chunk(
            asset.id,
            chunk_index=index,
            content=chunk.content,
            start_ms=chunk.start_ms,
            end_ms=chunk.end_ms,
            speaker=chunk.speaker,
            summary=chunk.summary,
        )
    return metadata


def _ingest_video(asset: SessionAssetSummary) -> dict[str, object]:
    original_path = Path(asset.storage_path)
    metadata = _probe_media_metadata(original_path, kind="video")
    transcript_status, transcript_chunks = _transcribe_video_asset(asset, metadata)
    metadata["transcript_status"] = transcript_status
    keyframe_status, keyframes = _extract_video_keyframes(asset, metadata)
    metadata["keyframe_status"] = keyframe_status
    if keyframes:
        metadata["keyframe_paths"] = [frame["path"] for frame in keyframes]
        metadata["keyframe_count"] = len(keyframes)
    for index, chunk in enumerate(transcript_chunks):
        asset_service.create_asset_chunk(
            asset.id,
            chunk_index=index,
            content=chunk.content,
            start_ms=chunk.start_ms,
            end_ms=chunk.end_ms,
            speaker=chunk.speaker,
            summary=chunk.summary,
        )
    for keyframe in keyframes:
        timestamp_ms = int(keyframe.get("timestamp_ms") or 0)
        frame_index = int(keyframe.get("frame_index") or 0)
        asset_service.create_asset_chunk(
            asset.id,
            chunk_index=100000 + frame_index,
            content=f"Video keyframe extracted at {timestamp_ms}ms.",
            frame_index=frame_index,
            frame_timestamp_ms=timestamp_ms,
            summary=f"Keyframe at {timestamp_ms}ms",
        )
    return metadata


def _transcribe_audio_asset(asset: SessionAssetSummary, metadata: dict[str, object]) -> tuple[str, list[ParsedChunk]]:
    if settings.jarvis_speech_recognition_provider != "volcengine":
        return "not_configured", []
    original_path = Path(asset.storage_path)
    audio_format = _detect_transcribable_audio_format(original_path, asset.mime_type)
    if not audio_format:
        return "unsupported_format", []
    try:
        transcript = speech_recognition_service.transcribe(
            SpeechRecognitionRequest(
                asset_id=asset.id,
                mime_type=asset.mime_type,
                path=original_path.as_posix(),
            )
        )
    except speech_recognition_service.SpeechRecognitionError:
        return "failed", []

    chunks: list[ParsedChunk] = []
    if transcript.segments:
        for segment in transcript.segments:
            summary_text = segment.text[:200] if segment.text else None
            chunks.append(
                ParsedChunk(
                    content=segment.text,
                    start_ms=segment.start_ms,
                    end_ms=segment.end_ms,
                    speaker=segment.speaker,
                    summary=summary_text,
                )
            )
    elif transcript.text:
        chunks.append(ParsedChunk(content=transcript.text, summary=transcript.text[:200]))
    metadata.update(transcript.metadata or {})
    if transcript.text:
        metadata["transcript_preview"] = transcript.text[:500]
    status = "ready" if (chunks or transcript.text) else "empty"
    return status, chunks


def _transcribe_video_asset(asset: SessionAssetSummary, metadata: dict[str, object]) -> tuple[str, list[ParsedChunk]]:
    if settings.jarvis_speech_recognition_provider != "volcengine":
        return "not_configured", []
    original_path = Path(asset.storage_path)
    audio_path = _extract_video_audio_track(asset.session_id, asset.id, original_path)
    if audio_path is None:
        return "no_audio", []
    try:
        transcript = speech_recognition_service.transcribe(
            SpeechRecognitionRequest(
                asset_id=asset.id,
                mime_type="audio/wav",
                path=audio_path.as_posix(),
            )
        )
    except speech_recognition_service.SpeechRecognitionError:
        return "failed", []

    chunks: list[ParsedChunk] = []
    if transcript.segments:
        for segment in transcript.segments:
            summary_text = segment.text[:200] if segment.text else None
            chunks.append(
                ParsedChunk(
                    content=segment.text,
                    start_ms=segment.start_ms,
                    end_ms=segment.end_ms,
                    speaker=segment.speaker,
                    summary=summary_text,
                )
            )
    elif transcript.text:
        chunks.append(ParsedChunk(content=transcript.text, summary=transcript.text[:200]))
    metadata.update(transcript.metadata or {})
    metadata["derived_audio_path"] = audio_path.as_posix()
    if transcript.text:
        metadata["transcript_preview"] = transcript.text[:500]
    status = "ready" if (chunks or transcript.text) else "empty"
    return status, chunks


def _extract_video_keyframes(asset: SessionAssetSummary, metadata: dict[str, object]) -> tuple[str, list[dict[str, object]]]:
    ffmpeg_path = _which_binary("ffmpeg")
    if not ffmpeg_path:
        return "not_configured", []
    video_path = Path(asset.storage_path)
    derived_dir = session_asset_utils.asset_derived_dir(asset.session_id, asset.id)
    derived_dir.mkdir(parents=True, exist_ok=True)
    duration_ms = metadata.get("duration_ms")
    if isinstance(duration_ms, int) and duration_ms > 0:
        timestamps_ms = _keyframe_timestamps(duration_ms, max_frames=3)
    else:
        timestamps_ms = [0]

    frames: list[dict[str, object]] = []
    for frame_index, timestamp_ms in enumerate(timestamps_ms):
        output_path = derived_dir / f"keyframe-{frame_index:03d}.png"
        command = [
            ffmpeg_path,
            "-y",
            "-ss",
            f"{max(0, timestamp_ms) / 1000:.3f}",
            "-i",
            video_path.as_posix(),
            "-frames:v",
            "1",
            output_path.as_posix(),
        ]
        try:
            subprocess.run(
                command,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.PIPE,
                text=True,
            )
        except subprocess.CalledProcessError:
            continue
        if output_path.exists() and output_path.stat().st_size > 0:
            frames.append(
                {
                    "frame_index": frame_index,
                    "timestamp_ms": timestamp_ms,
                    "path": output_path.as_posix(),
                }
            )
    return ("ready" if frames else "failed"), frames


def _extract_document_chunks(asset: SessionAssetSummary) -> list[ParsedChunk]:
    original_path = Path(asset.storage_path)
    if asset.kind == "pdf":
        return _extract_pdf_chunks(original_path)
    if asset.kind == "docx":
        return _extract_docx_chunks(original_path)
    if asset.kind == "xlsx":
        return _extract_xlsx_chunks(original_path)
    if asset.kind == "pptx":
        return _extract_pptx_chunks(original_path)
    raise AssetUploadError(f"Unsupported asset kind for ingestion: {asset.kind}")


def _probe_media_metadata(path: Path, *, kind: str) -> dict[str, object]:
    metadata: dict[str, object] = {}
    suffix = path.suffix.lower()
    header = b""
    try:
        with path.open("rb") as handle:
            header = handle.read(32)
    except Exception:
        header = b""

    if kind == "audio":
        audio_metadata = _probe_audio_metadata(path, suffix=suffix, header=header)
        metadata.update(audio_metadata)
        return metadata

    if kind == "video":
        video_metadata = _probe_video_metadata(path=path, suffix=suffix, header=header)
        metadata.update(video_metadata)
        return metadata

    return metadata


def _probe_audio_metadata(path: Path, *, suffix: str, header: bytes) -> dict[str, object]:
    metadata: dict[str, object] = {}
    container = ""
    if header.startswith(b"RIFF") and header[8:12] == b"WAVE":
        container = "wav"
    elif header.startswith(b"ID3"):
        container = "mp3"
    elif header.startswith(b"fLaC"):
        container = "flac"
    elif header.startswith(b"OggS"):
        container = "ogg"
    elif suffix == ".m4a":
        container = "m4a"
    elif suffix == ".aac":
        container = "aac"
    elif suffix == ".mp3":
        container = "mp3"
    if container:
        metadata["container"] = container

    if container == "wav":
        try:
            with wave.open(path.as_posix(), "rb") as handle:
                frame_rate = int(handle.getframerate() or 0)
                frame_count = int(handle.getnframes() or 0)
                channel_count = int(handle.getnchannels() or 0)
                sample_width = int(handle.getsampwidth() or 0)
                if frame_rate > 0 and frame_count >= 0:
                    metadata["duration_ms"] = int((frame_count / frame_rate) * 1000)
                    metadata["sample_rate"] = frame_rate
                if channel_count > 0:
                    metadata["channels"] = channel_count
                if sample_width > 0:
                    metadata["sample_width_bytes"] = sample_width
        except Exception:
            pass
    return metadata


def _detect_transcribable_audio_format(path: Path, mime_type: str) -> str | None:
    suffix = path.suffix.lower()
    if suffix in {".wav", ".mp3", ".ogg"}:
        return suffix
    lowered = (mime_type or "").lower()
    if lowered in {"audio/wav", "audio/x-wav", "audio/mpeg", "audio/ogg", "audio/opus"}:
        return lowered
    return None


def _probe_video_metadata(*, path: Path, suffix: str, header: bytes) -> dict[str, object]:
    metadata: dict[str, object] = {}
    if len(header) >= 12 and header[4:8] == b"ftyp":
        metadata["container"] = "mp4"
    elif header.startswith(b"\x1a\x45\xdf\xa3"):
        metadata["container"] = "matroska"
    elif suffix in {".mov", ".m4v"}:
        metadata["container"] = suffix.lstrip(".")
    elif suffix in {".mp4", ".webm", ".mkv"}:
        metadata["container"] = suffix.lstrip(".")
    metadata.update(_probe_video_stream_details(path))
    return metadata


def _probe_video_stream_details(path: Path) -> dict[str, object]:
    ffprobe_path = _which_binary("ffprobe")
    if not ffprobe_path:
        return {}
    command = [
        ffprobe_path,
        "-v",
        "error",
        "-print_format",
        "json",
        "-show_streams",
        "-show_format",
        path.as_posix(),
    ]
    try:
        result = subprocess.run(
            command,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
        )
        payload = json.loads(result.stdout)
    except Exception:
        return {}
    metadata: dict[str, object] = {}
    streams = payload.get("streams") if isinstance(payload, dict) else None
    if isinstance(streams, list):
        video_stream = next((stream for stream in streams if isinstance(stream, dict) and stream.get("codec_type") == "video"), None)
        audio_stream = next((stream for stream in streams if isinstance(stream, dict) and stream.get("codec_type") == "audio"), None)
        if isinstance(video_stream, dict):
            width = _as_int(video_stream.get("width"))
            height = _as_int(video_stream.get("height"))
            if width is not None:
                metadata["width"] = width
            if height is not None:
                metadata["height"] = height
        metadata["has_audio"] = audio_stream is not None
    format_info = payload.get("format") if isinstance(payload, dict) else None
    if isinstance(format_info, dict):
        duration_raw = format_info.get("duration")
        try:
            if duration_raw is not None:
                metadata["duration_ms"] = int(float(duration_raw) * 1000)
        except Exception:
            pass
    return metadata


def _extract_video_audio_track(session_id: str, asset_id: str, video_path: Path) -> Path | None:
    ffmpeg_path = _which_binary("ffmpeg")
    ffprobe_path = _which_binary("ffprobe")
    if not ffmpeg_path or not ffprobe_path:
        return None
    if not _video_has_audio_stream(ffprobe_path, video_path):
        return None
    derived_dir = session_asset_utils.asset_derived_dir(session_id, asset_id)
    derived_dir.mkdir(parents=True, exist_ok=True)
    output_path = derived_dir / "extracted-audio.wav"
    command = [
        ffmpeg_path,
        "-y",
        "-i",
        video_path.as_posix(),
        "-vn",
        "-acodec",
        "pcm_s16le",
        "-ac",
        "1",
        "-ar",
        "16000",
        output_path.as_posix(),
    ]
    try:
        subprocess.run(
            command,
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
            text=True,
        )
    except subprocess.CalledProcessError:
        return None
    return output_path if output_path.exists() and output_path.stat().st_size > 0 else None


def _video_has_audio_stream(ffprobe_path: str, video_path: Path) -> bool:
    command = [
        ffprobe_path,
        "-v",
        "error",
        "-select_streams",
        "a",
        "-show_entries",
        "stream=index",
        "-of",
        "csv=p=0",
        video_path.as_posix(),
    ]
    try:
        result = subprocess.run(
            command,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
        )
    except subprocess.CalledProcessError:
        return False
    return bool(result.stdout.strip())


def _which_binary(name: str) -> str | None:
    command = ["which", name]
    try:
        result = subprocess.run(
            command,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
        )
    except subprocess.CalledProcessError:
        return None
    resolved = result.stdout.strip()
    return resolved or None


def _keyframe_timestamps(duration_ms: int, *, max_frames: int) -> list[int]:
    if duration_ms <= 0 or max_frames <= 1:
        return [0]
    if duration_ms < 3000:
        return [0]
    positions = [0.2, 0.5, 0.8][:max_frames]
    timestamps = sorted({max(0, int(duration_ms * position)) for position in positions})
    return timestamps or [0]


def _chunk_text_blocks(
    blocks: list[tuple[str, dict[str, object]]],
    *,
    char_limit: int,
) -> list[ParsedChunk]:
    chunks: list[ParsedChunk] = []
    buffer: list[str] = []
    current_meta: dict[str, object] = {}

    def flush() -> None:
        if not buffer:
            return
        content = "\n".join(buffer).strip()
        if not content:
            buffer.clear()
            return
        chunks.append(
            ParsedChunk(
                content=content,
                page_number=_as_optional_int(current_meta.get("page_number")),
                sheet_name=_as_optional_str(current_meta.get("sheet_name")),
                slide_number=_as_optional_int(current_meta.get("slide_number")),
                section_path=_as_optional_str(current_meta.get("section_path")),
                start_ms=_as_optional_int(current_meta.get("start_ms")),
                end_ms=_as_optional_int(current_meta.get("end_ms")),
                speaker=_as_optional_str(current_meta.get("speaker")),
                frame_index=_as_optional_int(current_meta.get("frame_index")),
                frame_timestamp_ms=_as_optional_int(current_meta.get("frame_timestamp_ms")),
                summary=content[:200],
            )
        )
        buffer.clear()

    for text, meta in blocks:
        cleaned = text.strip()
        if not cleaned:
            continue
        meta_changed = buffer and meta != current_meta
        projected = len("\n".join(buffer + [cleaned]))
        if meta_changed or projected > char_limit:
            flush()
        if not buffer:
            current_meta = dict(meta)
        buffer.append(cleaned)
    flush()
    return chunks


def _as_optional_int(value: object) -> int | None:
    return value if isinstance(value, int) else None


def _as_optional_str(value: object) -> str | None:
    return value if isinstance(value, str) and value.strip() else None


def _extract_pdf_chunks(path: Path) -> list[ParsedChunk]:
    reader = PdfReader(path)
    blocks: list[tuple[str, dict[str, object]]] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        blocks.append((text, {"page_number": index}))
    return _chunk_text_blocks(blocks, char_limit=settings.jarvis_asset_chunk_char_limit)


def _extract_docx_chunks(path: Path) -> list[ParsedChunk]:
    document = DocxDocument(path)
    blocks: list[tuple[str, dict[str, object]]] = []
    current_heading = ""
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        if style_name.lower().startswith("heading"):
            current_heading = text
        blocks.append((text, {"section_path": current_heading or None}))
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            values = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if not values:
                continue
            blocks.append(
                ("\t".join(values), {"section_path": current_heading or f"table-{table_index}", "page_number": row_index})
            )
    return _chunk_text_blocks(blocks, char_limit=settings.jarvis_asset_chunk_char_limit)


def _extract_xlsx_chunks(path: Path) -> list[ParsedChunk]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    blocks: list[tuple[str, dict[str, object]]] = []
    try:
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if not values:
                    continue
                blocks.append(("\t".join(values), {"sheet_name": sheet.title}))
        return _chunk_text_blocks(blocks, char_limit=settings.jarvis_asset_chunk_char_limit)
    finally:
        workbook.close()


def _extract_pptx_chunks(path: Path) -> list[ParsedChunk]:
    presentation = Presentation(path)
    blocks: list[tuple[str, dict[str, object]]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        notes_text = ""
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame and notes_frame.text else ""
        except Exception:
            notes_text = ""
        if notes_text:
            texts.append(f"Notes: {notes_text}")
        if texts:
            blocks.append(("\n".join(texts), {"slide_number": index, "section_path": f"slide-{index}"}))
    return _chunk_text_blocks(blocks, char_limit=settings.jarvis_asset_chunk_char_limit)
